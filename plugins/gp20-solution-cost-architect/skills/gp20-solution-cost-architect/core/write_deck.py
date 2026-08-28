"""
Proposal deck — offering-agnostic.

Reads only the contract shape, so the same five slides serve managed LAN, DaaS
or anything else a pack produces. python-pptx only; no Node.

Usage:  python write_deck.py result.json [output.pptx]
"""

from __future__ import annotations

# --- Windows console safety -------------------------------------------------
import sys as _sys
for _s in (_sys.stdout, _sys.stderr):
    try:
        if _s and getattr(_s, "encoding", "").lower().replace("-", "") != "utf8":
            _s.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
# ---------------------------------------------------------------------------

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1F, 0x35, 0x56)
TEAL = RGBColor(0x2E, 0x6E, 0x8E)
MIST = RGBColor(0xF2, 0xF5, 0xF8)
CARD_LINE = RGBColor(0xDC, 0xE4, 0xEC)
AMBER = RGBColor(0xC8, 0x79, 0x1A)
AMBER_BG = RGBColor(0xFD, 0xF3, 0xE4)
INK = RGBColor(0x25, 0x32, 0x42)
MUTED = RGBColor(0x6B, 0x77, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COVER_SUB = RGBColor(0x8F, 0xB3, 0xC9)
COVER_LBL = RGBColor(0x7E, 0x9D, 0xB5)
COVER_TXT = RGBColor(0xCB, 0xD9, 0xE4)
GRID = RGBColor(0xE4, 0xEA, 0xF0)

HEAD, BODY = "Cambria", "Calibri"
W, H, M = 13.333, 7.5, 0.7


def _fmt(item: dict, symbol: str) -> str:
    v, f = item["value"], item.get("format", "int")
    if f == "money":
        return f"{symbol}{round(float(v)):,}"
    if f == "float1":
        return f"{float(v):,.1f}"
    if f == "float2":
        return f"{float(v):,.2f}"
    if f == "text":
        return str(v)
    return f"{round(float(v)):,}"


def money(n, symbol="£") -> str:
    """Sign goes outside the currency symbol: -£542,132, not £-542,132."""
    v = round(float(n))
    return f"-{symbol}{abs(v):,}" if v < 0 else f"{symbol}{v:,}"


def textbox(slide, x, y, w, h, text, *, size=13, bold=False, color=INK,
            font=BODY, align=PP_ALIGN.LEFT, italic=False, spacing=None,
            anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        if spacing is not None:
            p.space_after = Pt(spacing)
    return tb


def card(slide, x, y, w, h, *, fill=MIST, line=CARD_LINE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def stat(slide, x, y, w, h, value, label, *, hl=False, size=29):
    card(slide, x, y, w, h, fill=AMBER_BG if hl else MIST,
         line=AMBER if hl else CARD_LINE)
    textbox(slide, x, y + 0.16, w, h * 0.5, value, size=size, bold=True,
            color=AMBER if hl else NAVY, font=HEAD, align=PP_ALIGN.CENTER)
    textbox(slide, x, y + h * 0.62, w, h * 0.3, label, size=11,
            color=MUTED, align=PP_ALIGN.CENTER)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def heading(slide, kicker, title):
    textbox(slide, M, 0.42, W - 2 * M, 0.28, kicker.upper(), size=11,
            bold=True, color=TEAL)
    textbox(slide, M, 0.68, W - 2 * M, 0.75, title, size=31, bold=True,
            color=NAVY, font=HEAD)


def footer(slide, text):
    textbox(slide, M, 7.0, W - 2 * M, 0.28, text, size=8.5, italic=True, color=MUTED)


def style_chart(chart, colour, *, labels=False, fmt="0.0"):
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 60
    s = plot.series[0]
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = colour
    if labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = fmt
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        dl.font.size = Pt(10)
        dl.font.color.rgb = INK
        dl.font.name = BODY
    for axis, is_cat in ((chart.category_axis, True), (chart.value_axis, False)):
        axis.has_major_gridlines = not is_cat
        axis.major_tick_mark = XL_TICK_MARK.NONE
        axis.minor_tick_mark = XL_TICK_MARK.NONE
        axis.tick_labels.font.size = Pt(10.5)
        axis.tick_labels.font.name = BODY
        axis.tick_labels.font.color.rgb = INK if is_cat else MUTED
        if not is_cat:
            axis.major_gridlines.format.line.color.rgb = GRID
            axis.major_gridlines.format.line.width = Pt(0.75)
        axis.format.line.color.rgb = GRID


# ---------------------------------------------------------------------------

def _cover(prs, r):
    meta, summ = r["meta"], r["summary"]
    sym = meta["symbol"]
    s = blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            Emu(int(W * 914400)), Emu(int(H * 914400)))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    textbox(s, M, 2.25, W - 2 * M, 0.4, "SOLUTION & COST SUMMARY", size=13,
            bold=True, color=COVER_SUB)
    textbox(s, M, 2.7, W - 2 * M, 1.0, meta["client_name"], size=38, bold=True,
            color=WHITE, font=HEAD)
    textbox(s, M, 3.8, W - 2 * M - 3, 0.5, meta["offering_name"], size=17,
            color=COVER_TXT)

    facts = [("RFP reference", meta["rfp_ref"]),
             ("Contract term", f"{summ['term_years']} years")]
    if r["scope"]["headline"]:
        h0 = r["scope"]["headline"][0]
        facts.append((h0["label"], _fmt(h0, sym)))
    facts.append(("Total contract value", money(summ["tcv"], sym)))

    for i, (k, v) in enumerate(facts[:4]):
        x = M + i * 3.05
        textbox(s, x, 4.9, 2.9, 0.26, k.upper(), size=9, bold=True, color=COVER_LBL)
        textbox(s, x, 5.16, 2.9, 0.45, v, size=18, bold=True, color=WHITE, font=HEAD)

    textbox(s, M, 6.85, W - 2 * M, 0.3, meta["disclaimer"], size=9,
            italic=True, color=COVER_LBL)
    s.notes_slide.notes_text_frame.text = (
        "Open by stating this is an illustrative model on synthetic rates. The "
        "point is the workflow from RFP to costed solution, not the price.")


def _solution(prs, r):
    meta, dep = r["meta"], r["deployment"]
    sym = meta["symbol"]
    s = blank(prs)
    heading(s, "Solution shape", "What we would build")

    items = r["scope"]["headline"][:4]
    cw, gap = 2.86, 0.32
    for i, it in enumerate(items):
        stat(s, M + i * (cw + gap), 1.75, cw, 1.28, _fmt(it, sym), it["label"])

    textbox(s, M, 3.42, 5.6, 0.34, "Service commitment", size=13, bold=True, color=NAVY)
    for i, item in enumerate(r["service"][:6]):
        y = 3.86 + i * 0.38
        textbox(s, M, y, 2.9, 0.32, item["label"], size=12, color=MUTED)
        textbox(s, M + 2.9, y, 2.9, 0.32, str(item["value"]), size=12,
                bold=True, color=INK)

    card(s, 7.15, 3.42, 5.45, 2.72)
    textbox(s, 7.45, 3.62, 4.9, 0.32, "Transition effort", size=13, bold=True, color=NAVY)
    roles = [(k, v) for k, v in dep["days_by_role"].items() if k != "total"][:5]
    for i, (k, v) in enumerate(roles):
        y = 4.06 + i * 0.34
        textbox(s, 7.45, y, 3.4, 0.3, k.replace("_", " ").title(), size=11.5, color=MUTED)
        textbox(s, 10.85, y, 1.5, 0.3, f"{v:,.0f} days", size=11.5, bold=True,
                color=INK, align=PP_ALIGN.RIGHT)
    y = 4.06 + len(roles) * 0.34 + 0.12
    textbox(s, 7.45, y, 3.4, 0.32, "Total", size=12.5, bold=True, color=NAVY)
    textbox(s, 10.85, y, 1.5, 0.32, f"{dep['days_by_role'].get('total', 0):,.0f} days",
            size=12.5, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

    footer(s, meta["disclaimer"])


def _resource(prs, r):
    meta, run = r["meta"], r["run"]
    s = blank(prs)
    heading(s, "Resource model", "Who delivers it, and from where")

    by_loc: dict[str, float] = {}
    for res in run["resources"]:
        by_loc[res["location"]] = by_loc.get(res["location"], 0.0) + res["fte"]

    if by_loc:
        data = CategoryChartData()
        data.categories = list(by_loc)
        data.add_series("FTE", [round(v, 2) for v in by_loc.values()])
        gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M),
                                Inches(1.95), Inches(7.1), Inches(4.15), data)
        style_chart(gf.chart, TEAL, labels=True)
        textbox(s, M, 1.72, 7.1, 0.28, "Full-time equivalents by location", size=12,
                bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    stat(s, 8.2, 1.8, 4.4, 1.3, f"{run['total_fte']:,.1f}",
         "Total FTE in the managed service")

    has = bool(run.get("insight"))
    card(s, 8.2, 3.32, 4.4, 2.78, fill=AMBER_BG if has else MIST,
         line=AMBER if has else CARD_LINE)
    textbox(s, 8.48, 3.52, 3.9, 0.3, "Architect's observation", size=12,
            bold=True, color=AMBER if has else NAVY)
    textbox(s, 8.48, 3.86, 3.86, 2.05,
            run.get("insight") or "No structural observations for this configuration.",
            size=11, color=INK)

    footer(s, meta["disclaimer"])
    s.notes_slide.notes_text_frame.text = (
        "The slide to linger on. The observation is what separates an architect "
        "from a calculator — offer the sensitivity run here.")


def _commercials(prs, r):
    meta, dep, run, summ = r["meta"], r["deployment"], r["run"], r["summary"]
    sym = meta["symbol"]
    s = blank(prs)
    heading(s, "Commercials", "What it costs")

    cw, gap = 3.73, 0.355
    tiles = [(money(dep["price"], sym), "One-off charge", False),
             (money(run["price_pa"], sym), "Annual charge (year 1)", False),
             (money(summ["tcv"], sym),
              f"Total contract value · {summ['term_years']} years", True)]
    for i, (v, l, hl) in enumerate(tiles):
        stat(s, M + i * (cw + gap), 1.78, cw, 1.42, v, l, hl=hl)

    data = CategoryChartData()
    data.categories = [f"Year {y['year']}" for y in summ["yearly"]]
    data.add_series("Annual charge", [round(y["price"]) for y in summ["yearly"]])
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(3.72),
                            Inches(7.1), Inches(2.4), data)
    style_chart(gf.chart, NAVY)
    gf.chart.value_axis.tick_labels.number_format = f'"{sym}"#,##0,,"M"'
    gf.chart.value_axis.tick_labels.number_format_is_linked = False
    textbox(s, M, 3.5, 7.1, 0.28, "Recurring charge by contract year", size=12,
            bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    textbox(s, 8.2, 3.5, 4.4, 0.32, "Annual cost build", size=13, bold=True, color=NAVY)
    # Show every line if it fits; credits and contingency are the interesting
    # ones and always sit at the end.
    lines = run["cost_lines"][:9]
    step = 0.29 if len(lines) > 7 else 0.31
    for i, line in enumerate(lines):
        y = 3.9 + i * step
        textbox(s, 8.2, y, 2.5, 0.27, line["label"], size=10.5, color=MUTED)
        textbox(s, 10.7, y, 1.9, 0.27, money(line["amount"], sym), size=10.5,
                bold=True, color=AMBER if line["amount"] < 0 else INK,
                align=PP_ALIGN.RIGHT)

    if summ.get("unit_metrics"):
        u = summ["unit_metrics"][0]
        y = 3.9 + len(lines) * step + 0.12
        textbox(s, 8.2, y, 4.4, 0.3, f"{u['label']}: {sym}{u['value']:,.2f}",
                size=12, bold=True, color=NAVY)

    footer(s, meta["disclaimer"])


def _assumptions(prs, r):
    meta = r["meta"]
    s = blank(prs)
    heading(s, "Assumptions", "What we know, and what we need")

    flagged = {f["parameter"] for f in r.get("review_flags", [])}
    entries = r.get("assumptions", [])
    from_rfp = [a for a in entries if a["source"] == "rfp"]
    needs = [a for a in entries if a["source"] == "default" or a["parameter"] in flagged]

    def pretty(a) -> str:
        v = a["value"]
        if isinstance(v, dict):
            v = ", ".join(
                f"{k} {round(x * 100)}%" if isinstance(x, (int, float)) and x <= 1
                else f"{k} {x}" for k, x in v.items())
        elif a["parameter"].endswith("_pct") and isinstance(v, (int, float)):
            v = f"{v * 100:.1f}%"
        t = str(v)
        if len(t) > 42:
            t = t[:41].rstrip(", ") + "…"
        return f"•  {a['parameter'].replace('_', ' ')} — {t}"

    for x, items, title_, fill, line, hcol in [
        (M, from_rfp, f"Taken from the RFP  ·  {len(from_rfp)}", MIST, CARD_LINE, NAVY),
        (M + 6.15, needs, f"Requires confirmation  ·  {len(needs)}",
         AMBER_BG, AMBER, AMBER),
    ]:
        card(s, x, 1.8, 5.85, 4.35, fill=fill, line=line)
        textbox(s, x + 0.3, 2.02, 5.25, 0.34, title_, size=13, bold=True, color=hcol)
        body = ("\n".join(pretty(a) for a in items[:9]) if items
                else "•  Every parameter is traceable to the RFP or to you.")
        textbox(s, x + 0.3, 2.44, 5.25, 3.5, body, size=11.5, color=INK, spacing=5)

    footer(s, meta["disclaimer"])
    s.notes_slide.notes_text_frame.text = (
        "Close here. Nothing in the price is unattributed — every figure traces "
        "to the RFP, to a decision the user made, or to a flagged default.")


def build(result: dict, out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    for fn in (_cover, _solution, _resource, _commercials, _assumptions):
        fn(prs, result)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 1
    src = sys.argv[1]
    result = json.load(sys.stdin if src == "-" else open(src, encoding="utf-8"))
    out = sys.argv[2] if len(sys.argv) > 2 else "Proposal_Deck.pptx"
    print(f"Written: {build(result, out)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
